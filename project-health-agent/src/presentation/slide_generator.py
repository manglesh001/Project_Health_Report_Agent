from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.parser.data_model import Project
from src.scoring.rag_engine import RAGResult


RAG_COLORS = {
    "Red": RGBColor(0xE7, 0x4C, 0x3C),
    "Amber": RGBColor(0xF3, 0x9C, 0x12),
    "Green": RGBColor(0x27, 0xAE, 0x60),
}

DARK_BG = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)


def create_executive_deck(
    projects: list[tuple[Project, RAGResult]],
    synthesis_text: str,
    output_dir: str = "data/output",
) -> str:
    """generate a 5-7 slide executive presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # parse synthesis into sections
    sections = _parse_synthesis_sections(synthesis_text)

    def find_section(*keywords):
        """find a section by matching any keyword in the section key"""
        for key, val in sections.items():
            for kw in keywords:
                if kw in key:
                    return val
        return ""

    # slide 1: title
    _add_title_slide(prs, projects)

    # slide 2: portfolio overview
    _add_portfolio_slide(prs, projects, find_section("portfolio", "overview"))

    # slide 3: cross-project trends
    _add_text_slide(prs, "Cross-Project Trends", find_section("trend", "cross-project"))

    # slide 4: risk heat map
    _add_risk_slide(prs, projects, find_section("risk", "heat map"))

    # slide 5: top blockers
    _add_text_slide(prs, "Top Blockers & Dependencies", find_section("blocker", "dependencies"))

    # slide 6: delivery forecast + recommendations
    recs = find_section("recommend")
    forecast = find_section("forecast", "delivery")
    _add_text_slide(prs, "Delivery Forecast & Recommendations", f"{forecast}\n\n{recs}")

    # slide 7: next month outlook
    _add_text_slide(prs, "Next Month Outlook", find_section("outlook", "next month"))

    # save
    today = date.today()
    filename = f"monthly_executive_report_{today.strftime('%Y%m')}.pptx"
    path = Path(output_dir) / filename
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return str(path)


def _add_title_slide(prs: Presentation, projects: list[tuple[Project, RAGResult]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG

    # title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Monthly Project Health Report"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # subtitle
    p2 = tf.add_paragraph()
    p2.text = f"{date.today().strftime('%B %Y')} | {len(projects)} Active Projects"
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

    # RAG summary
    red = sum(1 for _, r in projects if r.overall_status == "Red")
    amber = sum(1 for _, r in projects if r.overall_status == "Amber")
    green = sum(1 for _, r in projects if r.overall_status == "Green")

    p3 = tf.add_paragraph()
    p3.text = f"\n🟢 {green}  🟡 {amber}  🔴 {red}"
    p3.font.size = Pt(24)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.LEFT


def _add_portfolio_slide(prs: Presentation, projects: list[tuple[Project, RAGResult]], text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # header
    _add_slide_header(slide, "Portfolio Overview")

    # project cards
    card_top = Inches(1.8)
    card_width = Inches(5.5)
    card_height = Inches(1.5)

    for i, (project, rag) in enumerate(projects):
        left = Inches(0.5) + (i % 2) * Inches(6.5)
        top = card_top + (i // 2) * Inches(1.8)

        shape = slide.shapes.add_shape(1, left, top, card_width, card_height)  # rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        color = RAG_COLORS.get(rag.overall_status, WHITE)
        p.text = f"● {project.summary.name}"
        p.font.size = Pt(14)
        p.font.bold = True

        p2 = tf.add_paragraph()
        p2.text = f"  {rag.overall_status} | {project.summary.percent_complete}% | {project.summary.current_stage}"
        p2.font.size = Pt(11)

        # top issue
        issues = [d.reason for d in rag.dimensions if d.status == "Red"]
        if not issues:
            issues = [d.reason for d in rag.dimensions if d.status == "Amber"]
        if issues:
            p3 = tf.add_paragraph()
            p3.text = f"  ⚠ {issues[0][:70]}"
            p3.font.size = Pt(10)

    # synthesis text below
    if text:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text[:400]
        p.font.size = Pt(11)


def _add_risk_slide(prs: Presentation, projects: list[tuple[Project, RAGResult]], text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, "Risk Heat Map")

    # risk table
    rows = 1 + len(projects)
    cols = 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5 + rows * 0.6))
    table = table_shape.table

    headers = ["Project", "Schedule", "Milestones", "Blockers", "Dependencies"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    for row_idx, (project, rag) in enumerate(projects):
        table.cell(row_idx + 1, 0).text = project.summary.name[:25]
        for col_idx, dim in enumerate(rag.dimensions):
            cell = table.cell(row_idx + 1, col_idx + 1)
            cell.text = dim.status
            # color the cell
            cell.fill.solid()
            cell.fill.fore_color.rgb = RAG_COLORS.get(dim.status, LIGHT_GRAY)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    # add text below table
    if text:
        top = Inches(2.5 + rows * 0.6)
        txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text[:500]
        p.font.size = Pt(11)


def _add_text_slide(prs: Presentation, title: str, text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, title)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    # split text into paragraphs
    paragraphs = text.strip().split("\n") if text else ["No data available."]
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # strip markdown headers
        clean = para_text.lstrip("#").strip()
        if para_text.startswith("- ") or para_text.startswith("* "):
            clean = "• " + clean[2:]
        p.text = clean
        p.font.size = Pt(12)
        if para_text.startswith("#") or para_text.startswith("**"):
            p.font.bold = True
            p.font.size = Pt(14)


def _add_slide_header(slide, title: str):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BG


def _parse_synthesis_sections(text: str) -> dict[str, str]:
    """split markdown synthesis into sections by ## headers"""
    import re
    sections = {}
    current_key = ""
    current_lines = []

    for line in text.split("\n"):
        if line.startswith("## "):
            if current_key:
                sections[current_key] = "\n".join(current_lines).strip()
            raw_key = line.lstrip("#").strip().lower()
            # strip leading numbers like "1. " or "1) "
            raw_key = re.sub(r"^\d+[\.\)]\s*", "", raw_key)
            current_key = raw_key
            current_lines = []
        else:
            current_lines.append(line)

    if current_key:
        sections[current_key] = "\n".join(current_lines).strip()

    return sections
